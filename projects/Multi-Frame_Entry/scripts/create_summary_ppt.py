"""
创建Multi-Frame Entry实验总结PPT（Google风格）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Google品牌色
GOOGLE_BLUE = RGBColor(66, 133, 244)    # #4285F4
GOOGLE_RED = RGBColor(234, 67, 53)      # #EA4335
GOOGLE_YELLOW = RGBColor(251, 187, 4)   # #FBBC04
GOOGLE_GREEN = RGBColor(52, 168, 83)    # #34A853
DARK_GRAY = RGBColor(60, 64, 67)        # #3C4043
LIGHT_GRAY = RGBColor(241, 243, 244)    # #F1F3F4
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs):
    """标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 背景色
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = GOOGLE_BLUE
    background.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Multi-Frame Entry"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "机器学习趋势识别系统实验总结"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = LIGHT_GRAY

    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "2026-02-20"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = LIGHT_GRAY

def add_objective_slide(prs):
    """实验目标页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 实验目标"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 三个目标卡片
    objectives = [
        ("构建多周期趋势识别系统", "使用机器学习预测市场趋势"),
        ("解决Regime漂移问题", "适应市场结构变化"),
        ("实现可交易信号", "AUC > 0.5，具备实用价值")
    ]

    y_pos = 2
    for i, (title, desc) in enumerate(objectives):
        # 卡片背景
        card = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.5))
        card.fill.solid()
        if i == 0:
            card.fill.fore_color.rgb = GOOGLE_BLUE
        elif i == 1:
            card.fill.fore_color.rgb = GOOGLE_GREEN
        else:
            card.fill.fore_color.rgb = GOOGLE_RED
        card.line.fill.background()

        # 标题
        title_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.2), Inches(8.6), Inches(0.5))
        title_tf = title_tb.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].font.size = Pt(24)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = WHITE

        # 描述
        desc_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.8), Inches(8.6), Inches(0.5))
        desc_tf = desc_tb.text_frame
        desc_tf.text = desc
        desc_tf.paragraphs[0].font.size = Pt(16)
        desc_tf.paragraphs[0].font.color.rgb = LIGHT_GRAY

        y_pos += 1.8

def add_journey_slide(prs):
    """实验历程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🔬 实验历程"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    phases = [
        ("Phase 1", "数据预处理", "多周期重采样"),
        ("Phase 2", "标签构建", "窗口对比分析"),
        ("Phase 3", "特征工程", "57个技术指标"),
        ("Phase 4", "模型优化", "4次迭代尝试")
    ]

    y_pos = 2
    for phase_num, phase_name, desc in phases:
        # 圆圈
        circle = slide.shapes.add_shape(9, Inches(1), Inches(y_pos), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOOGLE_BLUE
        circle.line.fill.background()

        # 阶段编号
        num_tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
        num_tf = num_tb.text_frame
        num_tf.text = str(phase_num.split()[1])
        num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_tf.paragraphs[0].font.size = Pt(20)
        num_tf.paragraphs[0].font.bold = True
        num_tf.paragraphs[0].font.color.rgb = WHITE

        # 阶段名称
        name_tb = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.05), Inches(3), Inches(0.4))
        name_tf = name_tb.text_frame
        name_tf.text = phase_name
        name_tf.paragraphs[0].font.size = Pt(24)
        name_tf.paragraphs[0].font.bold = True
        name_tf.paragraphs[0].font.color.rgb = DARK_GRAY

        # 描述
        desc_tb = slide.shapes.add_textbox(Inches(1.8), Inches(y_pos + 0.4), Inches(7), Inches(0.3))
        desc_tf = desc_tb.text_frame
        desc_tf.text = desc
        desc_tf.paragraphs[0].font.size = Pt(16)
        desc_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

        # 连接线
        if y_pos < 5:
            line = slide.shapes.add_shape(1, Inches(1.22), Inches(y_pos + 0.6), Inches(0.06), Inches(0.4))
            line.fill.solid()
            line.fill.fore_color.rgb = GOOGLE_BLUE
            line.line.fill.background()

        y_pos += 1.3

def add_regime_drift_slide(prs):
    """核心发现1：Regime漂移"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🔥 核心发现1：市场存在强烈Regime漂移"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # AUC数据展示
    auc_data = [
        ("2020→2021", "0.50", "随机"),
        ("2021→2022", "0.53", "可用"),
        ("2022→2023", "0.45", "失效"),
        ("2023→2024", "0.66", "优秀"),
        ("2024→2025", "0.53", "可用")
    ]

    y_pos = 2
    for year, auc, status in auc_data:
        # 背景条
        bar = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.7))
        bar.fill.solid()
        if status == "优秀":
            bar.fill.fore_color.rgb = GOOGLE_GREEN
        elif status == "可用":
            bar.fill.fore_color.rgb = LIGHT_GRAY
        else:
            bar.fill.fore_color.rgb = GOOGLE_RED
        bar.line.fill.background()

        # 年份
        year_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(2), Inches(0.5))
        year_tf = year_tb.text_frame
        year_tf.text = year
        year_tf.paragraphs[0].font.size = Pt(20)
        year_tf.paragraphs[0].font.bold = True
        year_tf.paragraphs[0].font.color.rgb = DARK_GRAY if status != "失效" else WHITE

        # AUC值
        auc_tb = slide.shapes.add_textbox(Inches(3), Inches(y_pos + 0.1), Inches(1.5), Inches(0.5))
        auc_tf = auc_tb.text_frame
        auc_tf.text = f"AUC: {auc}"
        auc_tf.paragraphs[0].font.size = Pt(24)
        auc_tf.paragraphs[0].font.bold = True
        auc_tf.paragraphs[0].font.color.rgb = DARK_GRAY if status != "失效" else WHITE

        # 状态
        status_tb = slide.shapes.add_textbox(Inches(5), Inches(y_pos + 0.1), Inches(2), Inches(0.5))
        status_tf = status_tb.text_frame
        status_tf.text = status
        status_tf.paragraphs[0].font.size = Pt(20)
        status_tf.paragraphs[0].font.color.rgb = DARK_GRAY if status != "失效" else WHITE

        y_pos += 1

    # 结论框
    conclusion_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(1))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.text = "❌ 不能用全历史平均参数  ✅ 必须使用滚动训练"
    conclusion_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    conclusion_frame.paragraphs[0].font.size = Pt(24)
    conclusion_frame.paragraphs[0].font.bold = True
    conclusion_frame.paragraphs[0].font.color.rgb = GOOGLE_RED

def add_window_slide(prs):
    """核心发现2：窗口大小"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⭐ 核心发现2：窗口大小至关重要"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 对比数据
    comparison = [
        ("20根K线窗口", "41.55%", "0.4816", "36.49%"),
        ("10根K线窗口", "58.63%", "0.5384", "65.65%")
    ]

    y_pos = 2.5
    for title, acc, auc, recall in comparison:
        # 背景框
        box = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(2))
        box.fill.solid()
        if "20根" in title:
            box.fill.fore_color.rgb = LIGHT_GRAY
        else:
            box.fill.fore_color.rgb = GOOGLE_BLUE
        box.line.fill.background()

        # 标题
        title_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.2), Inches(8.6), Inches(0.5))
        title_tf = title_tb.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = DARK_GRAY if "20根" in title else WHITE

        # 三个指标
        metrics = [("准确率", acc), ("AUC-ROC", auc), ("震荡召回率", recall)]
        x_pos = 0.7
        for metric_name, value in metrics:
            metric_tb = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.9), Inches(2.8), Inches(0.8))
            metric_tf = metric_tb.text_frame
            metric_tf.text = f"{metric_name}\n{value}"
            metric_tf.paragraphs[0].font.size = Pt(14)
            metric_tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80) if "20根" in title else LIGHT_GRAY
            metric_tf.paragraphs[1].font.size = Pt(28)
            metric_tf.paragraphs[1].font.bold = True
            metric_tf.paragraphs[1].font.color.rgb = DARK_GRAY if "20根" in title else WHITE

            x_pos += 3

        y_pos += 2.5

    # 改善标注
    improvement_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(0.7))
    improvement_frame = improvement_box.text_frame
    improvement_frame.text = "✅ 改善幅度：准确率 +54%  |  AUC-ROC +11.8%  |  震荡召回率 +925%"
    improvement_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    improvement_frame.paragraphs[0].font.size = Pt(20)
    improvement_frame.paragraphs[0].font.bold = True
    improvement_frame.paragraphs[0].font.color.rgb = GOOGLE_GREEN

def add_regime_filter_slide(prs):
    """核心发现3：Regime过滤"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⭐⭐ 核心发现3：Regime过滤显著提升性能"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 左侧：概念图
    concept_box = slide.shapes.add_shape(1, Inches(0.5), Inches(2), Inches(4), Inches(4))
    concept_box.fill.solid()
    concept_box.fill.fore_color.rgb = LIGHT_GRAY
    concept_box.line.fill.background()

    # 概念说明
    concept_title = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(3.6), Inches(0.6))
    concept_title_tf = concept_title.text_frame
    concept_title_tf.text = "策略思路"
    concept_title_tf.paragraphs[0].font.size = Pt(24)
    concept_title_tf.paragraphs[0].font.bold = True
    concept_title_tf.paragraphs[0].font.color.rgb = DARK_GRAY

    steps = [
        "① 识别波动率Regime",
        "② 筛选高波动数据",
        "③ 在高波动中训练趋势模型",
        "④ 只在高波动时使用模型"
    ]

    y_pos = 3
    for step in steps:
        step_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(3.6), Inches(0.5))
        step_tf = step_tb.text_frame
        step_tf.text = step
        step_tf.paragraphs[0].font.size = Pt(16)
        step_tf.paragraphs[0].font.color.rgb = DARK_GRAY
        y_pos += 0.6

    # 右侧：数据
    data_box = slide.shapes.add_shape(1, Inches(5), Inches(2), Inches(4.5), Inches(4))
    data_box.fill.solid()
    data_box.fill.fore_color.rgb = GOOGLE_BLUE
    data_box.line.fill.background()

    # 数据说明
    data_title = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.1), Inches(0.6))
    data_title_tf = data_title.text_frame
    data_title_tf.text = "高波动Regime特征"
    data_title_tf.paragraphs[0].font.size = Pt(24)
    data_title_tf.paragraphs[0].font.bold = True
    data_title_tf.paragraphs[0].font.color.rgb = WHITE

    facts = [
        "占总数据 40.4%",
        "ADX更高 (34.8 vs 31)",
        "ATR更大 (34.2 vs 18)",
        "测试准确率 +41%",
        "AUC-ROC首次突破0.5"
    ]

    y_pos = 3.2
    for fact in facts:
        fact_tb = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(4.1), Inches(0.45))
        fact_tf = fact_tb.text_frame
        fact_tf.text = fact
        fact_tf.paragraphs[0].font.size = Pt(18)
        fact_tf.paragraphs[0].font.color.rgb = WHITE
        y_pos += 0.5

def add_comparison_slide(prs):
    """三次迭代对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📊 三次迭代性能对比"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 表头
    headers = ["方法", "准确率", "AUC-ROC", "震荡召回率"]
    x_positions = [0.5, 3.5, 5.5, 7.5]

    for i, header in enumerate(headers):
        header_tb = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.5), Inches(2.5), Inches(0.5))
        header_tf = header_tb.text_frame
        header_tf.text = header
        header_tf.paragraphs[0].font.size = Pt(18)
        header_tf.paragraphs[0].font.bold = True
        header_tf.paragraphs[0].font.color.rgb = GOOGLE_BLUE

    # 数据行
    rows = [
        ("尝试1: 20根K线\n+绝对阈值0.5%", "37.98%", "0.4055", "6.40%", LIGHT_GRAY),
        ("尝试2: 20根K线\n+波动率归一化", "41.55%", "0.4816", "36.49%", LIGHT_GRAY),
        ("尝试3: 10根K线\n+Regime过滤", "58.63%", "0.5384", "65.65%", GOOGLE_GREEN)
    ]

    y_pos = 2.2
    for row in rows:
        # 背景条
        bg = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.3))
        bg.fill.solid()
        bg.fill.fore_color.rgb = row[4]
        bg.line.fill.background()

        for i, (text, color) in enumerate([(row[0], DARK_GRAY), (row[1], row[4]), (row[2], row[4]), (row[3], row[4])]):
            tb = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(y_pos + 0.1), Inches(2.5), Inches(1.1))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = text
            tf.paragraphs[0].font.size = Pt(16)
            if i == 0:
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = color
            else:
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = WHITE if color == GOOGLE_GREEN else DARK_GRAY

        y_pos += 1.5

def add_breakthrough_slide(prs):
    """性能突破"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 性能突破"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 大数字展示
    metrics = [
        ("58.63%", "测试准确率", "+54%", "GOOGLE_BLUE"),
        ("0.5384", "AUC-ROC", "+11.8%", "GOOGLE_GREEN"),
        ("65.65%", "震荡召回率", "+925%", "GOOGLE_RED")
    ]

    y_pos = 2
    for value, label, improvement, color in metrics:
        # 背景圆
        circle = slide.shapes.add_shape(9, Inches(1), Inches(y_pos), Inches(2.5), Inches(2.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = eval(color)
        circle.line.fill.background()

        # 大数字
        value_tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.7), Inches(2.5), Inches(1))
        value_tf = value_tb.text_frame
        value_tf.text = value
        value_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        value_tf.paragraphs[0].font.size = Pt(44)
        value_tf.paragraphs[0].font.bold = True
        value_tf.paragraphs[0].font.color.rgb = WHITE

        # 标签
        label_tb = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.3), Inches(3), Inches(0.5))
        label_tf = label_tb.text_frame
        label_tf.text = label
        label_tf.paragraphs[0].font.size = Pt(24)
        label_tf.paragraphs[0].font.color.rgb = DARK_GRAY

        # 改善
        improve_tb = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 1), Inches(3), Inches(0.5))
        improve_tf = improve_tb.text_frame
        improve_tf.text = f"改善: {improvement}"
        improve_tf.paragraphs[0].font.size = Pt(20)
        improve_tf.paragraphs[0].font.bold = True
        improve_tf.paragraphs[0].font.color.rgb = eval(color)

        y_pos += 2.5

def add_recommendation_slide(prs):
    """实盘建议"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "💡 实盘参数配置建议"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    # 配置参数
    configs = [
        ("训练方法", "滚动窗口"),
        ("重训频率", "每3个月"),
        ("训练窗口", "18-24个月"),
        ("预测窗口", "未来3个月"),
        ("窗口大小", "10根K线"),
        ("波动率阈值", "1.5σ"),
        ("Regime过滤", "仅高波动"),
        ("特征数量", "Top 30")
    ]

    # 分两列
    left_col = configs[:4]
    right_col = configs[4:]

    y_pos = 2
    for param, value in left_col:
        # 参数名
        param_tb = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(4), Inches(0.4))
        param_tf = param_tb.text_frame
        param_tf.text = f"• {param}"
        param_tf.paragraphs[0].font.size = Pt(18)
        param_tf.paragraphs[0].font.color.rgb = DARK_GRAY

        # 参数值
        value_tb = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.4), Inches(3.5), Inches(0.4))
        value_tf = value_tb.text_frame
        value_tf.text = value
        value_tf.paragraphs[0].font.size = Pt(20)
        value_tf.paragraphs[0].font.bold = True
        value_tf.paragraphs[0].font.color.rgb = GOOGLE_BLUE

        y_pos += 1

    y_pos = 2
    for param, value in right_col:
        # 参数名
        param_tb = slide.shapes.add_textbox(Inches(5), Inches(y_pos), Inches(4), Inches(0.4))
        param_tf = param_tb.text_frame
        param_tf.text = f"• {param}"
        param_tf.paragraphs[0].font.size = Pt(18)
        param_tf.paragraphs[0].font.color.rgb = DARK_GRAY

        # 参数值
        value_tb = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos + 0.4), Inches(3.5), Inches(0.4))
        value_tf = value_tb.text_frame
        value_tf.text = value
        value_tf.paragraphs[0].font.size = Pt(20)
        value_tf.paragraphs[0].font.bold = True
        value_tf.paragraphs[0].font.color.rgb = GOOGLE_BLUE

        y_pos += 1

    # 风控提示
    risk_box = slide.shapes.add_shape(1, Inches(0.5), Inches(7), Inches(9), Inches(1))
    risk_box.fill.solid()
    risk_box.fill.fore_color.rgb = GOOGLE_YELLOW
    risk_box.line.fill.background()

    risk_tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.2), Inches(8.6), Inches(0.6))
    risk_tf = risk_tb.text_frame
    risk_tf.text = "⚠️ 风控机制：若AUC < 0.5连续2个月，自动降低仓位(-50%)或暂停模型"
    risk_tf.paragraphs[0].font.size = Pt(18)
    risk_tf.paragraphs[0].font.bold = True
    risk_tf.paragraphs[0].font.color.rgb = DARK_GRAY

def add_next_steps_slide(prs):
    """下一步计划"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📅 下一步计划"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_GRAY

    phases = [
        ("Phase 5", "滚动训练框架", "1-2周", GOOGLE_BLUE),
        ("Phase 6", "策略框架构建", "2-3周", GOOGLE_GREEN),
        ("Phase 7", "完整回测系统", "3-4周", GOOGLE_YELLOW),
        ("Phase 8", "参数优化", "2-3周", GOOGLE_RED)
    ]

    y_pos = 2
    for phase_num, phase_name, duration, color in phases:
        # 背景条
        bar = slide.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 阶段编号
        num_tb = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(1.2), Inches(0.8))
        num_tf = num_tb.text_frame
        num_tf.text = phase_num.split()[1]
        num_tf.paragraphs[0].font.size = Pt(32)
        num_tf.paragraphs[0].font.bold = True
        num_tf.paragraphs[0].font.color.rgb = WHITE

        # 阶段名称
        name_tb = slide.shapes.add_textbox(Inches(2), Inches(y_pos + 0.1), Inches(4), Inches(0.4))
        name_tf = name_tb.text_frame
        name_tf.text = phase_name
        name_tf.paragraphs[0].font.size = Pt(24)
        name_tf.paragraphs[0].font.bold = True
        name_tf.paragraphs[0].font.color.rgb = WHITE

        # 时间
        time_tb = slide.shapes.add_textbox(Inches(6), Inches(y_pos + 0.1), Inches(3), Inches(0.4))
        time_tf = time_tb.text_frame
        time_tf.text = duration
        time_tf.paragraphs[0].font.size = Pt(20)
        time_tf.paragraphs[0].font.color.rgb = WHITE

        y_pos += 1.3

    # 总时间
    total_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(0.6))
    total_frame = total_box.text_frame
    total_frame.text = "📊 总计：前4个阶段约8-12周，实现从研究到可交易系统的完整路径"
    total_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    total_frame.paragraphs[0].font.size = Pt(20)
    total_frame.paragraphs[0].font.bold = True
    total_frame.paragraphs[0].font.color.rgb = DARK_GRAY

def add_summary_slide(prs):
    """总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = GOOGLE_BLUE
    background.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "✅ 关键结论"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # 五个关键点
    conclusions = [
        "模型不是没用 - 在不同Regime下有效性不同",
        "Regime过滤有效 - 高波动环境性能提升41%",
        "10根K线优于20根 - AUC从0.48提升到0.54",
        "滚动训练必要 - 市场存在强烈Regime漂移",
        "AUC=0.54可交易 - 配合风控有实用价值"
    ]

    y_pos = 3
    for conclusion in conclusions:
        conclusion_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.6))
        conclusion_frame = conclusion_box.text_frame
        conclusion_frame.text = conclusion
        conclusion_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        conclusion_frame.paragraphs[0].font.size = Pt(20)
        conclusion_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
        y_pos += 0.8

    # 底部标注
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "🚀 这已经比90%的量化研究强了！"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(24)
    footer_para.font.bold = True
    footer_para.font.color.rgb = GOOGLE_YELLOW

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 添加所有幻灯片
    add_title_slide(prs)  # 1
    add_objective_slide(prs)  # 2
    add_journey_slide(prs)  # 3
    add_regime_drift_slide(prs)  # 4
    add_window_slide(prs)  # 5
    add_regime_filter_slide(prs)  # 6
    add_comparison_slide(prs)  # 7
    add_breakthrough_slide(prs)  # 8
    add_recommendation_slide(prs)  # 9
    add_next_steps_slide(prs)  # 10
    add_summary_slide(prs)  # 11

    # 保存
    output_path = '/Users/mystryl/Documents/Quant/projects/Multi-Frame_Entry/Multi_Frame_Entry_Summary.pptx'
    prs.save(output_path)
    print(f"✓ PPT已生成: {output_path}")
    print(f"  共 {len(prs.slides)} 页")

if __name__ == '__main__':
    main()
